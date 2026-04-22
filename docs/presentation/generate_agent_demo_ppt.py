"""
generate_agent_demo_ppt.py
--------------------------
Generates a polished PPTX deck + flow diagram for the
Agentic AI Code Review Demo presentation.

Run from repo root:
    python docs/presentation/generate_agent_demo_ppt.py

Dependencies are installed automatically if missing.
Or install manually:
    pip install python-pptx Pillow
"""

import subprocess
import sys

# ── Auto-install missing dependencies ────────────────────────────────────────

def _ensure(*packages):
    for pkg in packages:
        try:
            __import__(pkg.replace("-", "_").split("[")[0])
        except ImportError:
            print(f"Installing {pkg} …")
            subprocess.check_call([sys.executable, "-m", "pip", "install", pkg])

_ensure("python-pptx", "Pillow")

# ── Imports (after install guard) ─────────────────────────────────────────────

from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Paths ────────────────────────────────────────────────────────────────────

ROOT       = Path(__file__).resolve().parent
PPTX_PATH  = ROOT / "Agentic_AI_Code_Review_Demo.pptx"
DIAG_PATH  = ROOT / "agentic-review-flow.png"

# ── Colour palette ────────────────────────────────────────────────────────────

PRIMARY = RGBColor(14,  44,  84)   # dark navy
ACCENT  = RGBColor(0,  129, 167)   # teal
SUCCESS = RGBColor(44, 120,  87)   # green
TEXT    = RGBColor(34,  34,  34)
MUTED   = RGBColor(100, 100, 100)
BG      = RGBColor(248, 250, 252)  # near-white page
LIGHT   = RGBColor(229, 238, 246)  # callout fill
WHITE   = RGBColor(255, 255, 255)

# ── Slide definitions ────────────────────────────────────────────────────────

SLIDES = [
    # 1 — Title
    {
        "type": "title",
        "title": "Agentic AI for Android PR Reviews",
        "subtitle": "Automated, context-aware code review — commented directly on GitHub",
        "footer":   "Built on Gemini 2.5 Flash · GitHub Actions · Python · MVVM Android project",
    },
    # 2 — Problem
    {
        "type": "bullets",
        "title": "The Problem",
        "bullets": [
            "Manual reviews are slow — developers wait hours or days for feedback.",
            "Quality varies by reviewer — same mistakes get through on different PRs.",
            "Repetitive checks waste senior engineer time: lint, magic numbers, hardcoded values.",
            "Test-coverage gaps go unnoticed until production.",
        ],
    },
    # 3 — Why Agentic
    {
        "type": "bullets",
        "title": "Why This Is an Agent — Not Just a Prompt",
        "bullets": [
            "Trigger-aware: starts automatically when a PR is opened or updated.",
            "Context-aware: reads PR diff, lint output, and test coverage.",
            "Role-aware: reasons through specialist reviewer personas (Architecture, UI, Security, QA).",
            "Action-oriented: posts structured review comments at the exact line of code.",
            "Feedback-loop: developer fixes → pushes → agent re-reviews in the same PR.",
        ],
        "callout": (
            "\"The agent behaves like a Staff Engineer: it reads context, "
            "reasons across specialties, and acts where the developer needs to.\""
        ),
    },
    # 4 — Flow diagram
    {
        "type": "diagram",
        "title": "End-to-End Flow",
        "caption": (
            "PR event → diff + lint + coverage collection → "
            "Gemini master reviewer → specialist sub-agents → "
            "structured findings → GitHub inline comments"
        ),
    },
    # 5 — Architecture 3-column
    {
        "type": "columns",
        "title": "How It Is Built",
        "columns": [
            {
                "heading": "① Orchestration",
                "items": [
                    "GitHub Actions workflow",
                    "Diff (pr_diff.patch)",
                    "Android Lint context",
                    "JaCoCo coverage XML",
                    "run_agent.py controller",
                ],
            },
            {
                "heading": "② Intelligence",
                "items": [
                    "master-reviewer-agent.md",
                    "arch-logic-agent.md",
                    "compose-ui-agent.md",
                    "security-config-agent.md",
                    "test-qa-agent.md",
                ],
            },
            {
                "heading": "③ Output",
                "items": [
                    "Structured XML findings",
                    "Severity ranking",
                    "Exact file + line pinning",
                    "Inline GitHub comments",
                    "Coverage alerts (< 75 %)",
                ],
            },
        ],
    },
    # 6 — What it checks
    {
        "type": "bullets",
        "title": "What the Agent Reviews",
        "bullets": [
            "Hardcoded strings, colours, magic numbers → suggests BuildConfig or resource XML.",
            "Android Lint + Kotlin lint rules and code formatting.",
            "Compose UI quality, accessibility, and recomposition hygiene.",
            "Architecture boundaries: ViewModel / Repository / Data layers.",
            "Security risks in build scripts, API usage, and Manifest.",
            "Test coverage gaps in business-logic files (Repository, UseCase).",
        ],
    },
    # 7 — Live demo
    {
        "type": "bullets",
        "title": "What the Live Demo Shows",
        "bullets": [
            "Open a PR with intentional issues in MainActivity.kt.",
            "GitHub Action auto-triggers: diff → lint → coverage → Gemini call.",
            "Agent returns structured XML findings ordered by severity.",
            "Each finding is posted as an inline comment on the exact changed line.",
            "Coverage alert posted separately when a business-logic file is below 75 %.",
        ],
        "callout": (
            "Demo tip: 4 bugs planted —\n"
            "hardcoded API key · hardcoded log tag\n"
            "magic number 5000 · hardcoded welcome string"
        ),
    },
    # 8 — Business impact
    {
        "type": "bullets",
        "title": "Business Impact",
        "bullets": [
            "Review feedback in ~2 minutes — not hours.",
            "Consistent quality: same bar on every PR, every repo.",
            "Senior engineer time freed from repetitive checks.",
            "Reusable policy: rules are version-controlled agent prompts.",
        ],
        "metrics": [
            ("Cycle time",   "Earlier feedback → fewer back-and-forth rounds"),
            ("Quality gate", "Lint + coverage + structured review in one automated pass"),
            ("Scale",        "Copy 3 files → works on any Android project"),
        ],
    },
    # 9 — How to reuse
    {
        "type": "bullets",
        "title": "How to Add This to Any Android Project",
        "bullets": [
            "Copy .github/workflows/ai-pr-reviewer.yml to your repo.",
            "Copy .github/scripts/run_agent.py and master-reviewer-agent.md.",
            "Add LLM_API_KEY to GitHub → Settings → Secrets.",
            "Optionally edit the .md prompt files to match your project's rules.",
            "Open a PR — the agent runs automatically.",
        ],
        "callout": "Total setup time: ~15 minutes",
    },
    # 10 — Close
    {
        "type": "bullets",
        "title": "Next Steps",
        "bullets": [
            "Pilot on one Android repo — run in review-only mode for 2 sprints.",
            "Tune the prompt files based on false-positive feedback.",
            "Track: review turnaround time, escaped defects, coverage trend.",
            "Templatise and roll out across additional mobile projects.",
        ],
        "callout": (
            "\"This is practical agentic AI: not generating text, "
            "but participating in an engineering workflow and producing "
            "accountable review actions.\""
        ),
    },
]

# ── PIL helpers ───────────────────────────────────────────────────────────────

def _pil_font(size: int) -> ImageFont.FreeTypeFont:
    for path in [
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Supplemental/Helvetica.ttc",
        "/System/Library/Fonts/Arial.ttf",
    ]:
        try:
            return ImageFont.truetype(path, size)
        except OSError:
            pass
    return ImageFont.load_default()

F_TITLE  = _pil_font(32)
F_BOX    = _pil_font(19)
F_CAPTION= _pil_font(16)


def _draw_box(draw, xy, text, fill, radius=16):
    """Draw a rounded rectangle with centred, auto-wrapped text."""
    x1, y1, x2, y2 = xy
    draw.rounded_rectangle(xy, radius=radius, fill=fill,
                           outline=(160, 180, 200), width=2)
    max_w = x2 - x1 - 20
    lines, current = [], ""
    for word in text.split():
        trial = f"{current} {word}".strip()
        if draw.textlength(trial, font=F_BOX) <= max_w:
            current = trial
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    lh = 23
    ty = y1 + (y2 - y1 - len(lines) * lh) / 2
    for ln in lines:
        tw = draw.textlength(ln, font=F_BOX)
        draw.text((x1 + (x2 - x1 - tw) / 2, ty), ln, font=F_BOX, fill=(20, 20, 20))
        ty += lh


def _arrow(draw, start, end):
    fill, w, head = (70, 100, 130), 4, 9
    draw.line([start, end], fill=fill, width=w)
    x1, y1 = start; x2, y2 = end
    if abs(x2 - x1) >= abs(y2 - y1):
        d = 1 if x2 >= x1 else -1
        draw.polygon([(x2, y2), (x2 - d*head*2, y2 - head), (x2 - d*head*2, y2 + head)], fill=fill)
    else:
        d = 1 if y2 >= y1 else -1
        draw.polygon([(x2, y2), (x2 - head, y2 - d*head*2), (x2 + head, y2 - d*head*2)], fill=fill)


# ── Flow diagram ──────────────────────────────────────────────────────────────

def create_flow_diagram(path: Path):
    W, H = 1600, 860
    img  = Image.new("RGB", (W, H), color=(248, 250, 252))
    draw = ImageDraw.Draw(img)

    # header
    draw.text((50, 28), "Agentic AI Review Flow", font=F_TITLE, fill=(14, 44, 84))
    draw.text((50, 72), "PR event → context collection → Gemini reasoning → GitHub comments", font=F_CAPTION, fill=(90, 90, 90))

    C_TRIGGER  = (220, 235, 252)
    C_CONTEXT  = (220, 245, 232)
    C_AGENT    = (255, 243, 210)
    C_OUTPUT   = (242, 220, 252)

    # ── Boxes (x1,y1,x2,y2) ─────────────────────────────────────────────────
    B = {
        # Column 1 — Trigger
        "pr":       (30,  160, 220, 240),
        # Column 2 — Orchestration
        "action":   (270, 160, 460, 240),
        "diff":     (510, 120, 700, 200),
        "lint":     (510, 240, 700, 320),
        "cov":      (510, 360, 700, 440),
        # Column 3 — Controller
        "runner":   (750, 230, 950, 330),
        # Column 4 — Master agent
        "master":   (1000, 230, 1200, 330),
        # Column 5 — Sub-agents
        "arch":     (1250, 140, 1450, 220),
        "ui":       (1250, 260, 1450, 340),
        "sec":      (1250, 380, 1450, 460),
        "qa":       (1250, 500, 1450, 580),
        # Column 6 — Output
        "findings": (1490, 240, 1570, 460),
        "comments": (750,  580, 1200, 660),
    }

    _draw_box(draw, B["pr"],       "Developer\nopens PR",         C_TRIGGER)
    _draw_box(draw, B["action"],   "GitHub Action\ntriggers",     C_TRIGGER)
    _draw_box(draw, B["diff"],     "PR diff\n(pr_diff.patch)",    C_CONTEXT)
    _draw_box(draw, B["lint"],     "Android Lint\noutput",        C_CONTEXT)
    _draw_box(draw, B["cov"],      "JaCoCo\ncoverage XML",        C_CONTEXT)
    _draw_box(draw, B["runner"],   "run_agent.py\norchestrator",  C_TRIGGER)
    _draw_box(draw, B["master"],   "master-reviewer\nagent",      C_AGENT)
    _draw_box(draw, B["arch"],     "Architecture\n& Logic",       C_AGENT)
    _draw_box(draw, B["ui"],       "Compose\nUI",                 C_AGENT)
    _draw_box(draw, B["sec"],      "Security\n& Config",          C_AGENT)
    _draw_box(draw, B["qa"],       "Test\n& QA",                  C_AGENT)
    _draw_box(draw, B["comments"], "GitHub inline review comments + coverage alerts", C_OUTPUT)

    def mid(box, axis="x"):
        return (box[0]+box[2])//2 if axis == "x" else (box[1]+box[3])//2

    # Trigger chain
    _arrow(draw, (B["pr"][2],    mid(B["pr"],"y")),   (B["action"][0], mid(B["action"],"y")))
    _arrow(draw, (B["action"][2],mid(B["action"],"y")),(B["diff"][0],   mid(B["diff"],"y")))
    _arrow(draw, (B["action"][2],mid(B["action"],"y")),(B["lint"][0],   mid(B["lint"],"y")))
    _arrow(draw, (B["action"][2],mid(B["action"],"y")),(B["cov"][0],    mid(B["cov"],"y")))

    # Context → runner
    _arrow(draw, (B["diff"][2],mid(B["diff"],"y")),   (B["runner"][0], mid(B["diff"],"y")))
    _arrow(draw, (B["lint"][2],mid(B["lint"],"y")),   (B["runner"][0], mid(B["lint"],"y")))
    _arrow(draw, (B["cov"][2], mid(B["cov"],"y")),    (B["runner"][0], mid(B["cov"],"y")))

    # runner → master → sub-agents
    _arrow(draw, (B["runner"][2],mid(B["runner"],"y")),(B["master"][0],mid(B["master"],"y")))
    for sub in ["arch","ui","sec","qa"]:
        _arrow(draw, (B["master"][2], mid(B["master"],"y")), (B[sub][0], mid(B[sub],"y")))

    # sub-agents → comments
    for sub in ["arch","ui","sec","qa"]:
        _arrow(draw, (mid(B[sub],"x"), B[sub][3]), (mid(B["comments"],"x"), B["comments"][1]))

    # Caption
    draw.text((50, 790), "Each finding is posted as an inline comment pinned to the exact changed line in the PR.", font=F_CAPTION, fill=(80, 80, 80))

    img.save(path)
    print(f"  Diagram  → {path}")


# ── PPTX helpers ──────────────────────────────────────────────────────────────

def _new_blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def _add_title(slide, text: str, y_in=0.32):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(y_in), Inches(12.2), Inches(0.75))
    p  = tb.text_frame.paragraphs[0]
    r  = p.add_run()
    r.text           = text
    r.font.name      = "Aptos"
    r.font.size      = Pt(26)
    r.font.bold      = True
    r.font.color.rgb = PRIMARY
    # underline rule
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(y_in + 0.72), Inches(12.2), Inches(0.03))
    rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()


def _add_callout(slide, text: str, x=8.0, y=1.5, w=4.9, h=None):
    lines = text.count("\n") + 1
    h = h or max(1.0, lines * 0.55 + 0.4)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = ACCENT
    tf = shape.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.text           = text
    p.font.name      = "Aptos"
    p.font.size      = Pt(15)
    p.font.bold      = True
    p.font.color.rgb = PRIMARY


def _add_bullets(slide, bullets, x=0.6, y=1.45, w=7.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text       = f"▸  {bullet}"
        p.space_after = Pt(10)
        p.font.name  = "Aptos"
        p.font.size  = Pt(19)
        p.font.color.rgb = TEXT


def _add_metrics(slide, metrics, x=7.9, start_y=3.8):
    for i, (label, detail) in enumerate(metrics):
        y = start_y + i * 0.82
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(y), Inches(5.0), Inches(0.68))
        shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = ACCENT
        tf = shape.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.text = f"  {label}:  {detail}"
        p.font.name  = "Aptos"
        p.font.size  = Pt(13)
        p.font.color.rgb = TEXT


# ── Slide builders ────────────────────────────────────────────────────────────

def _slide_title(prs, data):
    slide = _new_blank_slide(prs)
    # dark header band
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                  Inches(0), Inches(0), Inches(13.333), Inches(1.4))
    band.fill.solid(); band.fill.fore_color.rgb = PRIMARY
    band.line.fill.background()

    for shp in [
        slide.shapes.add_textbox(Inches(0.6), Inches(1.55), Inches(12), Inches(0.95)),
    ]:
        p = shp.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = data["title"]; r.font.name = "Aptos"
        r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = PRIMARY

    sub = slide.shapes.add_textbox(Inches(0.62), Inches(2.55), Inches(11), Inches(0.7))
    sp  = sub.text_frame.paragraphs[0]
    sp.text = data["subtitle"]
    sp.font.name = "Aptos"; sp.font.size = Pt(20); sp.font.color.rgb = TEXT

    highlight = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                       Inches(0.62), Inches(3.45), Inches(12.1), Inches(1.6))
    highlight.fill.solid(); highlight.fill.fore_color.rgb = LIGHT
    highlight.line.color.rgb = ACCENT
    htf = highlight.text_frame; htf.word_wrap = True
    hp  = htf.paragraphs[0]
    hp.text = (
        "Key idea: This is agentic AI — not just generating text, but participating in "
        "an engineering workflow and producing accountable, line-level review actions inside GitHub."
    )
    hp.font.name = "Aptos"; hp.font.size = Pt(19)
    hp.font.bold = True; hp.font.color.rgb = PRIMARY

    footer = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12), Inches(0.35))
    fp = footer.text_frame.paragraphs[0]
    fp.text = data["footer"]
    fp.font.name = "Aptos"; fp.font.size = Pt(11); fp.font.color.rgb = MUTED


def _slide_bullets(prs, data):
    slide = _new_blank_slide(prs)
    _add_title(slide, data["title"])
    w = 7.1 if ("callout" in data or "metrics" in data) else 12.3
    _add_bullets(slide, data["bullets"], w=w)
    if "callout" in data:
        _add_callout(slide, data["callout"])
    if "metrics" in data:
        _add_metrics(slide, data["metrics"])


def _slide_diagram(prs, data):
    slide = _new_blank_slide(prs)
    _add_title(slide, data["title"])
    slide.shapes.add_picture(str(DIAG_PATH), Inches(0.4), Inches(1.35), width=Inches(12.5))
    cap = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12), Inches(0.4))
    cp  = cap.text_frame.paragraphs[0]
    cp.text = data["caption"]
    cp.font.name = "Aptos"; cp.font.size = Pt(12); cp.font.color.rgb = MUTED


def _slide_columns(prs, data):
    slide = _new_blank_slide(prs)
    _add_title(slide, data["title"])
    positions = [0.5, 4.45, 8.4]
    for idx, col in enumerate(data["columns"]):
        x = Inches(positions[idx])
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                     x, Inches(1.45), Inches(3.5), Inches(5.3))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = ACCENT
        tf = box.text_frame; tf.word_wrap = True
        # heading
        ph = tf.paragraphs[0]
        ph.text = col["heading"]
        ph.font.name = "Aptos"; ph.font.size = Pt(18)
        ph.font.bold = True; ph.font.color.rgb = PRIMARY
        # items
        for item in col["items"]:
            pi = tf.add_paragraph()
            pi.text = f"  • {item}"
            pi.space_before = Pt(4)
            pi.font.name = "Aptos"; pi.font.size = Pt(14)
            pi.font.color.rgb = TEXT


# ── Build ────────────────────────────────────────────────────────────────────

BUILDERS = {
    "title":   _slide_title,
    "bullets": _slide_bullets,
    "diagram": _slide_diagram,
    "columns": _slide_columns,
}


def build_deck():
    ROOT.mkdir(parents=True, exist_ok=True)
    print("Generating flow diagram …")
    create_flow_diagram(DIAG_PATH)

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(SLIDES, 1):
        builder = BUILDERS[slide_data["type"]]
        builder(prs, slide_data)
        print(f"  Slide {i:>2}: {slide_data['title']}")

    prs.save(PPTX_PATH)
    print(f"\nDone → {PPTX_PATH}")


if __name__ == "__main__":
    build_deck()

